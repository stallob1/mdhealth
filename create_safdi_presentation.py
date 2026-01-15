#!/usr/bin/env python3
"""
Create a PowerPoint presentation about Dr. Alan V. Safdi's qualifications
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(0, 51, 102)  # Dark blue
    accent_color = RGBColor(0, 102, 204)  # Medium blue
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Dr. Alan V. Safdi, MD, FACG"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    subtitle.text = "CEO of MD Health\nWorld-Renowned Gastroenterologist & Wellness Advocate"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Executive Summary"
    title2.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content2.text = """• 47+ years of medical experience
• World-renowned gastroenterologist and researcher
• CEO of MD Health (nutrition bar company)
• Chief Medical Officer of Quadrant Health (Stanford/Mayo Clinic backed)
• 900+ clinical research studies involvement
• Serial entrepreneur with proven track record
• Passionate advocate for disease prevention and wellness"""
    
    for paragraph in content2.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
    
    # Slide 3: Education & Credentials
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Education & Credentials"
    title3.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content3.text = """UNDERGRADUATE
• Northwestern University (1971-1974)
• BA in Biology (graduated in 3 years)
• Phi Beta Kappa honors

MEDICAL SCHOOL
• University of Cincinnati College of Medicine (Class of 1978)
• Alpha Omega Alpha (AOA) Medical Honor Society

RESIDENCY & FELLOWSHIP
• UC San Diego Medical Center - Internal Medicine Residency
• University of Cincinnati - Gastroenterology Fellowship

BOARD CERTIFICATIONS
• Internal Medicine
• Gastroenterology
• Fellow of the American College of Gastroenterology (FACG)"""
    
    for paragraph in content3.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(8)
    
    # Slide 4: Current Leadership Positions
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Current Leadership Positions"
    title4.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content4.text = """CHIEF EXECUTIVE OFFICER
MD Health (nutrition bar company)

CHIEF MEDICAL OFFICER
Quadrant Health (co-owned by Stanford University and Mayo Clinic)

MEDICAL DIRECTOR
Telluride Longevity Institute

PRESIDENT
Consultants for Clinical Research"""
    
    for paragraph in content4.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(14)
    
    # Slide 5: Research & Publications
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Research & Publications"
    title5.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content5.text = """RESEARCH VOLUME
• Principal Investigator: 291 clinical research studies
• Co-Investigator: 845 clinical research studies
• Total: 900+ research protocols over 35+ years

PUBLICATION VENUES
• American Journal of Medicine
• American Journal of Gastroenterology
• Annals of Internal Medicine
• Alimentary Pharmacology & Therapeutics
• Numerous other prominent medical journals

RESEARCH FOCUS AREAS
• Nutrition and Gastroenterology
• Disease Prevention & Wellness
• Gut Health & Microbiome
• Longevity & Healthy Aging
• Inflammatory Bowel Disease (IBD)
• Colon Cancer Prevention"""
    
    for paragraph in content5.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(10)
    
    # Slide 6: Entrepreneurial Ventures
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Entrepreneurial Ventures"
    title6.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content6.text = """FOUNDED/CO-FOUNDED 11+ HEALTHCARE COMPANIES:

• MD Health - CEO (current nutrition bar company)
• Quadrant Health - CMO (AI-powered patient safety)
• eMerge Health Solutions - Co-founder & CEO
  (raised $850K from investors led by CincyTech)
• Ohio GI and Liver Institute - Co-founder
• Consultants for Clinical Research - Co-founder
• Telluride Longevity Institute - Co-founder
• Multiple Outpatient Endoscopy Centers
  (sold to national company within 3 years)
• CME Adventures, Veterinary Endo Institute, and more"""
    
    for paragraph in content6.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(10)
    
    # Slide 7: Media & Public Presence
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Media & Public Presence"
    title7.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content7.text = """PODCAST
• "Wellness Evidence-Based Medicine" (Apple Podcasts)
• Co-hosted with Dr. William Renner
• Evidence-based approach to wellness, nutrition, microbiome

REGULAR COLUMN
• "To Your Health" on Telluride Inside... and Out
• Topics: nutrition, longevity, disease prevention, supplements

SPEAKING ENGAGEMENTS
• Telluride Integrative Wellness Summit (with Deepak Chopra)
• Cornell University guest lectures
• "Unlocking Longevity" lecture at Sheridan Opera House
• International lecturer on wellness and gastroenterology

PROGRAMS CREATED
• "Live Longer Retreat" - Week-long wellness intensives in Telluride"""
    
    for paragraph in content7.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(10)
    
    # Slide 8: Philosophy & Approach
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Philosophy & Approach"
    title8.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content8.text = """CORE BELIEFS
"Passionate about disease prevention and wellness, 
not just fixing what has gone wrong."

KEY PRINCIPLES
✓ Evidence-Based Medicine
  All recommendations grounded in peer-reviewed research

✓ Prevention Over Treatment
  Focus on stopping disease before it starts

✓ Holistic Approach
  Incorporating diet, exercise, education, and alternative approaches

✓ Longevity Focus
  Studies populations with longest lifespans (Blue Zones)

✓ Mind-Body Connection
  Understanding relationship between mental and physical health"""
    
    for paragraph in content8.text_frame.paragraphs:
        paragraph.font.size = Pt(17)
        paragraph.space_after = Pt(10)
    
    # Slide 9: Awards & Recognition
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "Awards & Recognition"
    title9.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content9.text = """• Phi Beta Kappa - Prestigious academic honor society
• Alpha Omega Alpha (AOA) - Most prestigious medical honor society
  (elected in 3rd year of medical school)
• Fellow of the American College of Gastroenterology (FACG)
• Patient Top Choice Award - 5-star patient rating recognition
• President of Ohio Gastroenterology Society
  (elected to represent the State of Ohio)
• Advisory board member of Telluride First Foundation
• Featured on national program "Medical Crossfire"
• Principal lecturer alongside Deepak Chopra at wellness summits"""
    
    for paragraph in content9.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
    
    # Slide 10: Key Marketing Assets for MD Health
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Key Marketing Assets for MD Health"
    title10.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content10.text = """CREDIBILITY FACTORS
• 47+ years of medical experience
• 900+ research studies involvement
• Board certified in two specialties
• CMO of Quadrant Health (Stanford/Mayo Clinic backed)
• Publications in top medical journals
• Association with Deepak Chopra at wellness summits
• Cornell University guest lecturer
• Serial entrepreneur with proven track record

UNIQUE SELLING POINTS
• Doctor-formulated by a gastroenterologist (gut health expert)
• Backed by 35+ years of nutrition research
• Evidence-based approach to ingredients
• Created by someone who understands how food impacts 
  the digestive system
• Longevity and wellness focus aligns with health-conscious consumers"""
    
    for paragraph in content10.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(10)
    
    # Slide 11: Contact & Summary
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "Summary"
    title11.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content11.text = """Dr. Alan V. Safdi brings unparalleled credentials to MD Health:

• World-renowned gastroenterologist with 47+ years experience
• Extensive research background (900+ studies)
• Proven entrepreneurial success (11+ companies founded)
• Evidence-based approach to nutrition and wellness
• Passion for prevention over treatment
• Strong public presence and media credibility

His expertise in gut health, nutrition research, and longevity 
makes him the ideal authority figure for MD Health's mission 
of providing evidence-based nutrition products that promote 
wellness and disease prevention."""
    
    for paragraph in content11.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
    
    # Save presentation
    output_file = "Dr_Alan_Safdi_Qualifications.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
