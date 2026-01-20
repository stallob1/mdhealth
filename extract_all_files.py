#!/usr/bin/env python3
"""
Universal File Extractor
Extracts text from PDFs, PPTX, DOCX, XLSX, and other file types
"""

import os
import sys
from pathlib import Path

def extract_pdf(pdf_path):
    """Extract text from PDF"""
    try:
        import pdfplumber
        text = []
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text.append(f"\n--- Page {page_num + 1} ---\n")
                text.append(page.extract_text() or "")
        return '\n'.join(text)
    except Exception as e:
        return f"Error extracting PDF: {e}"

def extract_pptx(pptx_path):
    """Extract text from PowerPoint"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides):
            text.append(f"\n--- Slide {slide_num + 1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(shape.text)
                # Check for tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text.append(" | ".join(row_text))
        
        return '\n'.join(text) if text else "No text found in presentation"
    except Exception as e:
        return f"Error extracting PPTX: {e}"

def extract_docx(docx_path):
    """Extract text from Word document"""
    try:
        from docx import Document
        doc = Document(docx_path)
        text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            text.append("\n--- Table ---\n")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text.append(" | ".join(row_text))
        
        return '\n'.join(text) if text else "No text found in document"
    except Exception as e:
        return f"Error extracting DOCX: {e}"

def extract_xlsx(xlsx_path):
    """Extract text from Excel file"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path)
        text = []
        
        for sheet_name in wb.sheetnames:
            text.append(f"\n--- Sheet: {sheet_name} ---\n")
            ws = wb[sheet_name]
            
            for row in ws.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if any(row_text):
                    text.append(" | ".join(row_text))
        
        return '\n'.join(text) if text else "No data found in spreadsheet"
    except Exception as e:
        return f"Error extracting XLSX: {e}"

def extract_audio(audio_path):
    """Extract text from audio file (requires transcription)"""
    try:
        import speech_recognition as sr
        from pydub import AudioSegment
        
        # Convert audio to WAV if needed
        audio = AudioSegment.from_file(audio_path)
        wav_path = audio_path.with_suffix('.wav')
        audio.export(wav_path, format="wav")
        
        # Transcribe
        r = sr.Recognizer()
        with sr.AudioFile(str(wav_path)) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
        
        # Clean up temp file
        wav_path.unlink()
        return text
    except Exception as e:
        return f"Error transcribing audio: {e}\nNote: Audio transcription requires internet connection and may have limitations."

def extract_image(image_path):
    """Extract text from image using OCR"""
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text if text.strip() else "No text found in image"
    except Exception as e:
        return f"Error extracting text from image: {e}\nNote: Requires Tesseract OCR to be installed: brew install tesseract"

def extract_file(file_path):
    """Extract text from any supported file type"""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    
    extractors = {
        '.pdf': extract_pdf,
        '.pptx': extract_pptx,
        '.docx': extract_docx,
        '.xlsx': extract_xlsx,
        '.mp3': extract_audio,
        '.m4a': extract_audio,
        '.wav': extract_audio,
        '.png': extract_image,
        '.jpg': extract_image,
        '.jpeg': extract_image,
        '.gif': extract_image,
        '.tiff': extract_image,
    }
    
    if suffix not in extractors:
        return None
    
    print(f"Extracting text from {file_path.name}...")
    return extractors[suffix](file_path)

def main():
    """Main function"""
    if len(sys.argv) < 2:
        print("Usage: python3 extract_all_files.py <file> [output_file]")
        print("   or: python3 extract_all_files.py --all [output_folder]")
        print("\nSupported formats:")
        print("  PDF, PPTX, DOCX, XLSX, MP3, M4A, WAV, PNG, JPG, GIF, TIFF")
        sys.exit(1)
    
    if sys.argv[1] == '--all':
        output_folder = Path(sys.argv[2]) if len(sys.argv) > 2 else Path('extracted_text')
        output_folder.mkdir(exist_ok=True)
        
        # Find all supported files
        extensions = ['.pdf', '.pptx', '.docx', '.xlsx', '.mp3', '.m4a', '.wav', 
                     '.png', '.jpg', '.jpeg', '.gif', '.tiff']
        files = []
        for ext in extensions:
            files.extend(Path('.').glob(f'*{ext}'))
            files.extend(Path('.').glob(f'*{ext.upper()}'))
        
        if not files:
            print("No supported files found")
            return
        
        print(f"Found {len(files)} file(s). Extracting...\n")
        success_count = 0
        
        for file_path in files:
            text = extract_file(file_path)
            if text:
                output_file = output_folder / f"{file_path.stem}_extracted.txt"
                try:
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(f"Extracted from: {file_path.name}\n")
                        f.write("=" * 80 + "\n\n")
                        f.write(text)
                    print(f"✓ {file_path.name} -> {output_file.name}")
                    success_count += 1
                except Exception as e:
                    print(f"✗ Error writing {file_path.name}: {e}")
            else:
                print(f"✗ Unsupported format: {file_path.name}")
        
        print(f"\n✓ Successfully extracted {success_count}/{len(files)} file(s)")
    else:
        # Extract single file
        file_path = Path(sys.argv[1])
        output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else file_path.with_suffix('.txt')
        
        text = extract_file(file_path)
        if text:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"Extracted from: {file_path.name}\n")
                f.write("=" * 80 + "\n\n")
                f.write(text)
            print(f"✓ Extracted: {file_path.name} -> {output_path.name}")
        else:
            print(f"✗ Unsupported file format: {file_path.name}")

if __name__ == '__main__':
    main()
